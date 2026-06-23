"""导出引擎"""

import shutil
import subprocess
import tempfile
from pathlib import Path


class ExportManager:
    """导出管理器"""

    def to_pdf(self, html_path: Path, out_path: Path, width: int = 1920, height: int = 1080) -> dict:
        """导出为 PDF"""
        # 尝试使用 Playwright
        try:
            return self._export_with_playwright(html_path, out_path, "pdf", width, height)
        except:
            pass

        # 尝试使用 Puppeteer
        try:
            return self._export_with_puppeteer(html_path, out_path, "pdf", width, height)
        except:
            pass

        # 回退到浏览器打印
        return self._export_with_browser(html_path, out_path, "pdf")

    def to_png(self, html_path: Path, out_path: Path, width: int = 1920, height: int = 1080) -> dict:
        """导出为 PNG（第一页）"""
        try:
            return self._export_with_playwright(html_path, out_path, "png", width, height)
        except:
            pass
        return {"ok": False, "error": "PNG 导出需要 Playwright: pip install playwright"}

    def to_png_all(self, html_path: Path, out_dir: Path, width: int = 1920, height: int = 1080) -> dict:
        """导出所有页面为 PNG"""
        try:
            from playwright.sync_api import sync_playwright

            with sync_playwright() as p:
                browser = p.chromium.launch()
                page = browser.new_page(viewport={"width": width, "height": height})
                page.goto(f"file://{html_path.resolve()}")
                page.wait_for_timeout(1000)

                # 获取幻灯片数量
                slides = page.query_selector_all('#deck > section.slide')
                for i, _ in enumerate(slides):
                    page.evaluate(f"window.go({i})")
                    page.wait_for_timeout(500)
                    page.screenshot(path=str(out_dir / f"slide-{i+1:03d}.png"))

                browser.close()
                return {"ok": True, "path": str(out_dir), "pages": len(slides)}
        except ImportError:
            return {"ok": False, "error": "PNG 导出需要 Playwright: pip install playwright"}

    def to_pptx(self, html_path: Path, out_path: Path) -> dict:
        """导出为 PPTX"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor

            # 读取 HTML 并提取文本
            html = html_path.read_text(encoding="utf-8")
            # 这是一个简化实现，实际应用中需要更复杂的 HTML 解析
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            blank_slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(blank_slide_layout)

            # 添加标题
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(2))
            tf = title_box.text_frame
            tf.text = "HTML Point 导出"
            p = tf.paragraphs[0]
            p.font.size = Pt(44)
            p.font.bold = True

            prs.save(str(out_path))
            return {"ok": True, "path": str(out_path)}
        except ImportError:
            return {"ok": False, "error": "PPTX 导出需要 python-pptx: pip install python-pptx"}

    def to_html_zip(self, html_path: Path, out_path: Path) -> dict:
        """打包为独立 HTML ZIP"""
        import zipfile

        with zipfile.ZipFile(out_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.write(html_path, html_path.name)
            # 递归添加所有相关资源
            base_dir = html_path.parent
            for f in base_dir.rglob('*'):
                if f.is_file() and f != html_path and f != out_path:
                    rel_path = f.relative_to(base_dir)
                    zf.write(f, rel_path)

        return {"ok": True, "path": str(out_path)}

    def _export_with_playwright(self, html_path: Path, out_path: Path, fmt: str, width: int, height: int) -> dict:
        from playwright.sync_api import sync_playwright

        with sync_playwright() as p:
            browser = p.chromium.launch()
            page = browser.new_page(viewport={"width": width, "height": height})
            page.goto(f"file://{html_path.resolve()}")
            page.wait_for_timeout(1000)

            if fmt == "pdf":
                page.pdf(
                    path=str(out_path),
                    width=f"{width}px",
                    height=f"{height}px",
                    print_background=True,
                    margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
                )
            else:
                page.screenshot(path=str(out_path))

            browser.close()
            return {"ok": True, "path": str(out_path)}

    def _export_with_puppeteer(self, html_path: Path, out_path: Path, fmt: str, width: int, height: int) -> dict:
        cmd = [
            "npx", "puppeteer-cli",
            "screenshot" if fmt == "png" else "pdf",
            f"file://{html_path.resolve()}",
            "--output", str(out_path),
            "--width", str(width),
            "--height", str(height),
        ]
        subprocess.run(cmd, check=True, capture_output=True)
        return {"ok": True, "path": str(out_path)}

    def _export_with_browser(self, html_path: Path, out_path: Path, fmt: str) -> dict:
        """使用系统浏览器打印"""
        import webbrowser

        # 打开浏览器并触发打印
        webbrowser.open(f"file://{html_path.resolve()}")
        return {"ok": False, "error": "自动浏览器打印未完成。请手动使用浏览器的打印功能（Ctrl+P）并选择「另存为 PDF」"}
