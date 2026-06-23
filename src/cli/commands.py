# CLI 命令实现
import json
import os
import subprocess
import sys
import time
from pathlib import Path

from cli.utils import (
    print_error, print_success, print_info, print_warning,
    print_table, ensure_dir, load_json_or_yaml
)
from engine.builder import SlideBuilder
from engine.parser import parse_markdown, parse_json
from engine.theme import ThemeManager
from engine.layout import LayoutManager
from engine.export import ExportManager


def cmd_init(args) -> dict:
    """初始化新演示稿"""
    name = args.name.replace(" ", "-").lower()
    out_path = Path(args.out) if args.out else Path(f"{name}.html")
    theme = args.theme
    layout = args.layout
    num_slides = args.slides

    print_info(f"初始化演示稿: {name}")
    print_info(f"主题: {theme} · 布局: {layout} · 幻灯片: {num_slides}")

    builder = SlideBuilder(theme=theme)
    # 创建基础结构
    slides = []
    # 封面
    slides.append(builder.build_slide(layout="title", data={
        "title": name.replace("-", " ").title(),
        "subtitle": "由 HTML Point 生成",
    }))
    # 内容页
    for i in range(1, num_slides - 1):
        slides.append(builder.build_slide(layout="title-content", data={
            "title": f"第 {i} 页",
            "content": "在此处编辑内容...",
        }))
    # 结尾页
    slides.append(builder.build_slide(layout="title", data={
        "title": "谢谢",
        "subtitle": "Questions?",
    }))

    html = builder.assemble(name.replace("-", " ").title(), slides)
    out_path.write_text(html, encoding="utf-8")
    print_success(f"已创建: {out_path}")
    return {"ok": True, "path": str(out_path)}


def cmd_build(args) -> dict:
    """从数据文件构建 HTML 幻灯片"""
    input_path = Path(args.input)
    out_path = Path(args.out)
    theme = args.theme

    if not input_path.exists():
        return {"ok": False, "error": f"输入文件不存在: {input_path}"}

    print_info(f"构建: {input_path} → {out_path}")
    print_info(f"主题: {theme}")

    ext = input_path.suffix.lower()
    if ext == ".md":
        data = parse_markdown(input_path.read_text(encoding="utf-8"))
    elif ext in (".json", ".yaml", ".yml"):
        data = load_json_or_yaml(input_path)
    else:
        return {"ok": False, "error": f"不支持的输入格式: {ext}"}

    builder = SlideBuilder(theme=theme)
    slides = []
    for slide_data in data.get("slides", []):
        layout = slide_data.get("layout", "title-content")
        slides.append(builder.build_slide(layout=layout, data=slide_data))

    title = data.get("title", "Untitled")
    html = builder.assemble(title, slides)
    out_path.write_text(html, encoding="utf-8")
    print_success(f"已构建: {out_path}")
    return {"ok": True, "path": str(out_path)}


def cmd_serve(args) -> dict:
    """启动编辑服务器"""
    port = args.port
    open_browser = args.open and not args.no_browser

    print_info(f"启动服务器: http://localhost:{port}")
    from server import main as server_main
    import socketserver
    from src.app import Handler, PRESENTATIONS
    from src.storage import copy_demo
    from pathlib import Path

    # 创建 presentations 目录
    PRESENTATIONS.mkdir(exist_ok=True)
    # 复制 demo（如果存在）
    try:
        copy_demo(Path(__file__).resolve().parent.parent.parent / "2026机器人入职各行各业-演示.html", PRESENTATIONS)
    except:
        pass

    socketserver.TCPServer.allow_reuse_address = True
    with socketserver.TCPServer(("127.0.0.1", port), Handler) as httpd:
        print(f"\n  ◈  HTML Point v3 已启动")
        print(f"  📂 文稿: {PRESENTATIONS}")
        print(f"  🌐 http://localhost:{port}")
        print(f"  ⌨️  Ctrl+C 停止\n")
        if open_browser:
            subprocess.Popen(["open", f"http://localhost:{port}"])
        try:
            httpd.serve_forever()
        except KeyboardInterrupt:
            print("\n  已停止")
    return {"ok": True}


def cmd_convert(args) -> dict:
    """转换外部格式"""
    input_path = Path(args.input)
    out_dir = Path(args.out) if args.out else Path(".")
    theme = args.theme

    if not input_path.exists():
        return {"ok": False, "error": f"输入文件不存在: {input_path}"}

    ext = input_path.suffix.lower()
    if ext == ".pptx":
        return _convert_pptx(input_path, out_dir, theme)
    elif ext in (".md", ".markdown"):
        return _convert_md(input_path, out_dir, theme)
    else:
        return {"ok": False, "error": f"不支持的转换格式: {ext}"}


def _convert_pptx(input_path: Path, out_dir: Path, theme: str) -> dict:
    """将 PPTX 转换为 HTML"""
    try:
        from pptx import Presentation
    except ImportError:
        return {"ok": False, "error": "转换 PPTX 需要 python-pptx: pip install python-pptx"}

    print_info(f"转换 PPTX: {input_path}")
    prs = Presentation(str(input_path))
    builder = SlideBuilder(theme=theme)
    slides = []

    for slide in prs.slides:
        # 提取文本
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)

        title = texts[0] if texts else "Untitled"
        content = "\n".join(texts[1:]) if len(texts) > 1 else ""
        slides.append(builder.build_slide(layout="title-content", data={
            "title": title,
            "content": content,
        }))

    out_path = out_dir / f"{input_path.stem}.html"
    html = builder.assemble(input_path.stem, slides)
    out_path.write_text(html, encoding="utf-8")
    print_success(f"已转换: {out_path}")
    return {"ok": True, "path": str(out_path)}


def _convert_md(input_path: Path, out_dir: Path, theme: str) -> dict:
    """将 Markdown 转换为 HTML"""
    print_info(f"转换 Markdown: {input_path}")
    text = input_path.read_text(encoding="utf-8")
    data = parse_markdown(text)

    builder = SlideBuilder(theme=theme)
    slides = []
    for slide_data in data.get("slides", []):
        layout = slide_data.get("layout", "title-content")
        slides.append(builder.build_slide(layout=layout, data=slide_data))

    out_path = out_dir / f"{input_path.stem}.html"
    html = builder.assemble(data.get("title", input_path.stem), slides)
    out_path.write_text(html, encoding="utf-8")
    print_success(f"已转换: {out_path}")
    return {"ok": True, "path": str(out_path)}


def cmd_export(args) -> dict:
    """导出演示稿"""
    input_path = Path(args.input)
    fmt = args.format
    out_path = Path(args.out) if args.out else input_path.with_suffix(f".{fmt}")
    width = args.width
    height = args.height

    if not input_path.exists():
        return {"ok": False, "error": f"输入文件不存在: {input_path}"}

    print_info(f"导出: {input_path} → {out_path} (格式: {fmt})")
    exporter = ExportManager()

    if fmt == "pdf":
        result = exporter.to_pdf(input_path, out_path, width=width, height=height)
    elif fmt == "png":
        if args.all_pages:
            out_dir = out_path.parent / out_path.stem
            out_dir.mkdir(parents=True, exist_ok=True)
            result = exporter.to_png_all(input_path, out_dir, width=width, height=height)
        else:
            result = exporter.to_png(input_path, out_path, width=width, height=height)
    elif fmt == "pptx":
        result = exporter.to_pptx(input_path, out_path)
    elif fmt == "html-zip":
        result = exporter.to_html_zip(input_path, out_path)
    else:
        return {"ok": False, "error": f"不支持的导出格式: {fmt}"}

    if result.get("ok"):
        print_success(f"已导出: {result.get('path', out_path)}")
    return result


def cmd_theme(args) -> dict:
    """主题管理"""
    tm = ThemeManager()

    if args.theme_cmd == "list" or not args.theme_cmd:
        themes = tm.list_themes()
        print_table(
            ["名称", "描述", "颜色"],
            [[t["key"], t["name"], ", ".join(t.get("colors", [])[:3])] for t in themes]
        )
        return {"ok": True}

    elif args.theme_cmd == "apply":
        theme_key = args.theme
        file_path = Path(args.file)
        if not file_path.exists():
            return {"ok": False, "error": f"文件不存在: {file_path}"}

        out_path = Path(args.out) if args.out else file_path
        html = file_path.read_text(encoding="utf-8")
        new_html = tm.apply_theme(html, theme_key)
        out_path.write_text(new_html, encoding="utf-8")
        print_success(f"已应用主题 {theme_key} 到: {out_path}")
        return {"ok": True, "path": str(out_path)}

    return {"ok": False, "error": f"未知子命令: {args.theme_cmd}"}


def cmd_template(args) -> dict:
    """布局模板管理"""
    lm = LayoutManager()

    if args.template_cmd == "list" or not args.template_cmd:
        layouts = lm.list_layouts()
        print_table(
            ["名称", "描述", "标签"],
            [[l["name"], l.get("description", ""), ", ".join(l.get("tags", []))] for l in layouts]
        )
        return {"ok": True}

    elif args.template_cmd == "info":
        info = lm.get_layout_info(args.name)
        if not info:
            return {"ok": False, "error": f"未知布局: {args.name}"}
        print(json.dumps(info, ensure_ascii=False, indent=2))
        return {"ok": True}

    return {"ok": False, "error": f"未知子命令: {args.template_cmd}"}


def cmd_validate(args) -> dict:
    """验证演示稿格式"""
    file_path = Path(args.file)
    if not file_path.exists():
        return {"ok": False, "error": f"文件不存在: {file_path}"}

    html = file_path.read_text(encoding="utf-8")
    errors = []
    warnings = []

    # 基本结构检查
    if "<html" not in html.lower():
        errors.append("缺少 <html> 标签")
    if "<head" not in html.lower():
        errors.append("缺少 <head> 标签")
    if "<body" not in html.lower():
        errors.append("缺少 <body> 标签")

    # 幻灯片结构检查
    if '#deck' not in html and 'deck' not in html:
        warnings.append("未检测到 #deck 容器")
    if 'section' not in html or '.slide' not in html:
        warnings.append("未检测到 .slide 幻灯片")

    # 安全警告
    if '<script' in html.lower() and 'ppt-editor' not in html and 'ppt-presenter' not in html:
        warnings.append("包含外部 script 标签（可能的安全风险）")

    result = {
        "ok": len(errors) == 0,
        "file": str(file_path),
        "errors": errors,
        "warnings": warnings,
        "valid": len(errors) == 0,
    }

    if args.json:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        if result["valid"]:
            print_success("验证通过 ✓")
        else:
            print_error(f"验证失败: {len(errors)} 个错误")
        for e in errors:
            print_error(f"  - {e}")
        for w in warnings:
            print_warning(f"  - {w}")

    return result
