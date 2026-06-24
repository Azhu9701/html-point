"""PPTX 导出模块

将 HTML Point 演示稿导出为 .pptx 文件。
使用 Playwright 截图每张幻灯片, 嵌入到 PPTX 页面中。
"""

import io
import subprocess
import time
import socket
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def _find_free_port():
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


SERVER_SCRIPT = r"""
import sys, os, time
from pathlib import Path
from http.server import HTTPServer, BaseHTTPRequestHandler

project = Path(r'{PROJECT}')
html_file = Path(r'{HTML_FILE}')

sys.path.insert(0, str(project))
from src import security, editor_injector

html = html_file.read_text(encoding='utf-8')
html = security.strip_scripts_whitelist(html)

# 注入 presenter + core
core_js = project / 'web' / 'core.js'
pres_js = project / 'web' / 'presenter.js'
if core_js.exists():
    html = editor_injector.inject(html, str(core_js), marker='ppt-core')
if pres_js.exists():
    html = editor_injector.inject(html, str(pres_js), marker='ppt-presenter')

# 隐藏 UI
html = html.replace('</head>',
    '<style>#ppt-bar,#ppt-sb,#ppt-in,#ppt-gd,#ppt-zm,#ppt-to,#ppt-notes-panel,'
    '#ppt-pres-bar,.pc,.dh,.ppt-page-controls,.ppt-drag-handle,#nav,#hint'
    '{{display:none!important}}'
    ' body.pe .slide{{padding-left:5vw!important;padding-right:5vw!important}}'
    ' body.pe .canvas-card{{margin-left:0!important;width:100vw!important}}</style></head>', 1)

class H(BaseHTTPRequestHandler):
    def do_GET(self):
        self.send_response(200)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.end_headers()
        self.wfile.write(html.encode('utf-8'))
    def log_message(self, *a): pass

s = HTTPServer(('127.0.0.1', {PORT}), H)
s.serve_forever()
"""


def export_to_pptx(html_file: Path, output_path: Path) -> Path:
    """将 HTML 演示稿转为 PPTX (截图式, 每页一张全屏图片)"""
    from playwright.sync_api import sync_playwright

    html_file = html_file.resolve()
    output_path = output_path.resolve()
    project_dir = html_file.parent.parent.resolve()
    port = _find_free_port()

    script = SERVER_SCRIPT.format(
        PROJECT=str(project_dir),
        HTML_FILE=str(html_file),
        PORT=port,
    )

    server_proc = subprocess.Popen(
        [sys.executable, "-c", script],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )

    for _ in range(30):
        try:
            with socket.socket() as s:
                s.connect(("127.0.0.1", port))
                break
        except (ConnectionRefusedError, OSError):
            time.sleep(0.3)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1920, "height": 1080})
            page.goto(f"http://127.0.0.1:{port}/", wait_until="domcontentloaded", timeout=30000)
            time.sleep(3)

            slide_count = page.evaluate(
                "document.querySelectorAll('#deck > section.slide').length"
            )
            print(f"   检测到 {slide_count} 张幻灯片")

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for i in range(slide_count):
                page.evaluate(f"window.go({i})")
                time.sleep(1.0)

                screenshot = page.screenshot(full_page=False)
                img_stream = io.BytesIO(screenshot)

                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.add_picture(
                    img_stream,
                    left=0, top=0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
                print(f"   [{i+1}/{slide_count}] ✓")

            prs.save(str(output_path))
            browser.close()
    finally:
        server_proc.kill()
        server_proc.wait()

    return output_path
